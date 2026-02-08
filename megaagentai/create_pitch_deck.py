from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import nsmap

# Create presentation with 16:9 aspect ratio
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors
DARK_BG = RGBColor(15, 15, 25)
ACCENT_PURPLE = RGBColor(138, 43, 226)
ACCENT_BLUE = RGBColor(59, 130, 246)
WHITE = RGBColor(255, 255, 255)
LIGHT_GRAY = RGBColor(200, 200, 210)
DARK_GRAY = RGBColor(40, 40, 50)

def add_dark_background(slide):
    """Add dark background to slide"""
    background = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height
    )
    background.fill.solid()
    background.fill.fore_color.rgb = DARK_BG
    background.line.fill.background()
    # Send to back
    spTree = slide.shapes._spTree
    sp = background._element
    spTree.remove(sp)
    spTree.insert(2, sp)

def add_text_box(slide, left, top, width, height, text, font_size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    """Helper to add styled text box"""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return txBox

# ============================================
# SLIDE 1: Introduction - Meet the Founder
# ============================================
slide1_layout = prs.slide_layouts[6]  # Blank layout
slide1 = prs.slides.add_slide(slide1_layout)
add_dark_background(slide1)

# Title
add_text_box(slide1, 0.5, 0.3, 12.3, 0.8, "MegaAgentAI", font_size=24, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)

# Main heading
add_text_box(slide1, 0.5, 1.0, 12.3, 0.8, "Meet the Founder", font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Founder name
add_text_box(slide1, 0.5, 2.0, 12.3, 0.6, "Chandra Suda", font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Role
add_text_box(slide1, 0.5, 2.7, 12.3, 0.5, "Founder & CEO, MegaAgentAI", font_size=24, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Divider line
divider = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5), Inches(3.4), Inches(3.3), Inches(0.02))
divider.fill.solid()
divider.fill.fore_color.rgb = ACCENT_PURPLE
divider.line.fill.background()

# Education section
add_text_box(slide1, 0.5, 3.7, 12.3, 0.4, "🎓  EDUCATION", font_size=16, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide1, 0.5, 4.1, 12.3, 0.5, "B.S., Stanford University", font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Experience section  
add_text_box(slide1, 0.5, 4.9, 12.3, 0.4, "💼  EXPERIENCE", font_size=16, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide1, 0.5, 5.3, 12.3, 0.5, "Member of Technical Staff at xAI", font_size=32, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text_box(slide1, 0.5, 5.9, 12.3, 0.4, "Building AI systems at the frontier of intelligence", font_size=18, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Credential boxes
stanford_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4), Inches(6.5), Inches(2.3), Inches(0.6))
stanford_box.fill.solid()
stanford_box.fill.fore_color.rgb = DARK_GRAY
stanford_box.line.fill.background()
stanford_tf = stanford_box.text_frame
stanford_tf.paragraphs[0].text = "Stanford"
stanford_tf.paragraphs[0].font.size = Pt(18)
stanford_tf.paragraphs[0].font.bold = True
stanford_tf.paragraphs[0].font.color.rgb = RGBColor(140, 21, 21)
stanford_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

xai_box = slide1.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7), Inches(6.5), Inches(2.3), Inches(0.6))
xai_box.fill.solid()
xai_box.fill.fore_color.rgb = DARK_GRAY
xai_box.line.fill.background()
xai_tf = xai_box.text_frame
xai_tf.paragraphs[0].text = "xAI"
xai_tf.paragraphs[0].font.size = Pt(18)
xai_tf.paragraphs[0].font.bold = True
xai_tf.paragraphs[0].font.color.rgb = WHITE
xai_tf.paragraphs[0].alignment = PP_ALIGN.CENTER

# ============================================
# SLIDE 2: Problem, Value & Differentiation
# ============================================
slide2_layout = prs.slide_layouts[6]
slide2 = prs.slides.add_slide(slide2_layout)
add_dark_background(slide2)

# Title
add_text_box(slide2, 0.5, 0.2, 12.3, 0.5, "MegaAgentAI for Sama Club", font_size=20, bold=True, color=ACCENT_BLUE, align=PP_ALIGN.CENTER)
add_text_box(slide2, 0.5, 0.6, 12.3, 0.6, "Problem · Value · Differentiation", font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Left column - Problem
add_text_box(slide2, 0.4, 1.4, 4, 0.4, "🚨 THE PROBLEM", font_size=16, bold=True, color=RGBColor(239, 68, 68))

problems = [
    "Feedback scattered across channels",
    "12-16 weeks to get insights",
    "Reactive—damage done first",
    "Sentiment only, no root cause"
]
y_pos = 1.85
for problem in problems:
    add_text_box(slide2, 0.5, y_pos, 4, 0.35, f"• {problem}", font_size=15, color=LIGHT_GRAY)
    y_pos += 0.38

# Middle column - Business Value
add_text_box(slide2, 4.6, 1.4, 4, 0.4, "💰 BUSINESS VALUE", font_size=16, bold=True, color=RGBColor(34, 197, 94))

values = [
    ("⚡ Time to Insight", "24-72 hours"),
    ("🎯 Root-Cause", "Know WHY, not just WHAT"),
    ("📉 Reduce Churn", "Act before they leave"),
    ("🔮 Predictive", "Forecast issues early"),
]
y_pos = 1.85
for label, value in values:
    add_text_box(slide2, 4.7, y_pos, 4, 0.22, label, font_size=14, bold=True, color=WHITE)
    add_text_box(slide2, 4.7, y_pos + 0.22, 4, 0.22, value, font_size=13, color=LIGHT_GRAY)
    y_pos += 0.48

# Right column - Differentiation  
add_text_box(slide2, 8.8, 1.4, 4, 0.4, "🏆 WHY WE WIN", font_size=16, bold=True, color=ACCENT_PURPLE)

diffs = [
    ("Others: Dashboards", "Us: Decisions"),
    ("Others: Text only", "Us: Voice+Text+Images"),
    ("Others: Sentiment", "Us: Root Causes"),
    ("Others: Batch", "Us: Real-Time"),
]
y_pos = 1.85
for other, us in diffs:
    add_text_box(slide2, 8.9, y_pos, 4, 0.22, other, font_size=12, color=RGBColor(150, 150, 160))
    add_text_box(slide2, 8.9, y_pos + 0.22, 4, 0.22, f"→ {us}", font_size=13, bold=True, color=RGBColor(34, 197, 94))
    y_pos += 0.48

# Divider
divider2 = slide2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(4.1), Inches(12.3), Inches(0.015))
divider2.fill.solid()
divider2.fill.fore_color.rgb = DARK_GRAY
divider2.line.fill.background()

# Multi-Agent Architecture section
add_text_box(slide2, 0.5, 4.3, 12.3, 0.4, "🤖 Our 6 Specialized AI Agents", font_size=20, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

agents = [
    ("🎯", "Voice-of-Customer", "Tone & emotion"),
    ("🔄", "Normalizer", "Unify channels"),
    ("🛡️", "Authenticity", "Filter spam/bots"),
    ("🔍", "Investigator", "Find root causes"),
    ("🧠", "Learning Layer", "Gets smarter"),
    ("📊", "Orchestrator", "Coordinate all"),
]

x_pos = 0.6
for emoji, name, desc in agents:
    agent_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x_pos), Inches(4.8), Inches(1.95), Inches(1.1))
    agent_box.fill.solid()
    agent_box.fill.fore_color.rgb = DARK_GRAY
    agent_box.line.color.rgb = ACCENT_PURPLE
    
    add_text_box(slide2, x_pos + 0.05, 4.85, 1.85, 0.35, f"{emoji} {name}", font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide2, x_pos + 0.05, 5.2, 1.85, 0.35, desc, font_size=11, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    x_pos += 2.1

# Key metric callout
metric_box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(3.5), Inches(6.15), Inches(6.3), Inches(0.7))
metric_box.fill.solid()
metric_box.fill.fore_color.rgb = ACCENT_PURPLE
metric_box.line.fill.background()

add_text_box(slide2, 3.5, 6.25, 6.3, 0.5, "Time to Insight: 12-16 weeks → 24-72 hours", font_size=22, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Tagline
add_text_box(slide2, 0.5, 7.0, 12.3, 0.35, '"Most tools show what happened. We tell you why—and what to do next."', font_size=14, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)

# Save the presentation
output_path = "/Users/akshu/megaAI/megaagentai/MegaAgentAI_SamaClub_Pitch.pptx"
prs.save(output_path)
print(f"✅ Presentation saved to: {output_path}")
